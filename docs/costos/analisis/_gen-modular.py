"""
Genera ExamLab-Presentacion-Modelo-Modular.pptx a partir del modelo modular v3
(docs/costos/analisis/modelo-modular-v3.md): SIN plan Free, base por matrículas
+ add-ons a la carta, con bundles coherentes por perfil.

Uso:  py docs/costos/analisis/_gen-modular.py
Requiere: python-pptx. Misma paleta indigo/violeta (16:9, Calibri) que los decks
originales (Comercial/Aliados/General/etc.).
Salida → docs/demos/presentacion/ExamLab-Presentacion-Modelo-Modular.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Paleta compartida con los decks originales (indigo/violeta), NO teal.
# Idéntica a _gen-presentacion.py y a Comercial/Aliados/General/etc.
TEAL_DARK  = RGBColor(0x1E, 0x1B, 0x4B)   # Titular — indigo muy oscuro
TEAL       = RGBColor(0x4F, 0x46, 0xE5)   # Acento primario — indigo
TEAL_LIGHT = RGBColor(0x7C, 0x3A, 0xED)   # Acento claro — violeta
BLUE_ACCENT= RGBColor(0x4F, 0x46, 0xE5)   # Azul/indigo (igual al original; era #1E5BAF)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT    = RGBColor(0xF5, 0xF3, 0xFF)   # Fondo suave — violeta-50 (original)
TEXT_MAIN  = RGBColor(0x33, 0x41, 0x55)   # Texto principal — slate (original)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)   # Texto atenuado — slate (original)
GOLD       = RGBColor(0xE1, 0x9A, 0x1F)
GREEN_OK   = RGBColor(0x05, 0x96, 0x69)   # Verde (original)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_text_box(slide, x, y, w, h, text="", *, font="Calibri", size=14, bold=False,
                 italic=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None, line_width_pt=1.0, no_line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.05
    if fill_rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    if no_line:
        shape.line.fill.background()
    elif line_rgb is not None:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    shape.shadow.inherit = False
    return shape


def add_footer(slide, n):
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
                 "ExamLab · Modelo modular — planes + add-ons", size=9, color=TEXT_MUTED)
    add_text_box(slide, Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3),
                 f"{n}", size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def header(slide, kicker, title):
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.4),
                 kicker.upper(), size=12, bold=True, color=TEAL_LIGHT)
    add_text_box(slide, Inches(0.5), Inches(0.78), Inches(12.3), Inches(0.8),
                 title, size=30, bold=True, color=TEAL_DARK)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s_portada(prs):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=TEAL_DARK, no_line=True)
    add_text_box(s, Inches(1), Inches(2.4), Inches(11.3), Inches(0.5),
                 "EXAMLAB", size=16, bold=True, color=TEAL_LIGHT)
    add_text_box(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2),
                 "Modelo modular", size=54, bold=True, color=BG_WHITE)
    add_text_box(s, Inches(1), Inches(4.4), Inches(11.3), Inches(0.6),
                 "Un plan base por tamaño + add-ons a la carta. Sin versión gratuita: pagas por lo que usas.",
                 size=18, color=RGBColor(0xC7, 0xD2, 0xFE))


def s_filosofia(prs):
    s = blank(prs)
    header(s, "Cómo se arma", "Dos capas: un plan base + los add-ons que necesites")
    # Base card
    add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(4.4),
             fill_rgb=RGBColor(0xF5, 0xF3, 0xFF), line_rgb=TEAL, line_width_pt=2.0)
    add_text_box(s, Inches(0.8), Inches(2.15), Inches(5.4), Inches(0.5),
                 "1 · Plan base (obligatorio)", size=18, bold=True, color=TEAL_DARK)
    add_text_box(s, Inches(0.8), Inches(2.75), Inches(5.4), Inches(0.5),
                 "Según tus MATRÍCULAS activas. Define escala, storage y soporte.",
                 size=13, color=TEXT_MAIN)
    for i, t in enumerate(["Pequeña — hasta 1.000", "Mediana — hasta 3.000",
                            "Grande — hasta 10.000", "Enterprise — a medida",
                            "(Starter $79 solo para pilotos ≤200)"]):
        add_text_box(s, Inches(1.0), Inches(3.45 + i * 0.5), Inches(5.2), Inches(0.4),
                     ("• " + t), size=13, color=TEXT_MAIN if i < 4 else TEXT_MUTED,
                     italic=(i == 4))
    # Add-ons card
    add_rect(s, Inches(6.83), Inches(1.9), Inches(6.0), Inches(4.4),
             fill_rgb=BG_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=1.0)
    add_text_box(s, Inches(7.13), Inches(2.15), Inches(5.4), Inches(0.5),
                 "2 · Add-ons (opcionales, a la carta)", size=18, bold=True, color=TEAL_DARK)
    for i, t in enumerate(["IA administrada — $0,10/matrícula",
                           "Storage extra — $10 / 100 GB",
                           "Code runner (Java/Python) — $49",
                           "Aislamiento dedicado — $99",
                           "SSO / SAML — $29 (+$99 setup)",
                           "Certificación oficial — $29"]):
        add_text_box(s, Inches(7.33), Inches(2.85 + i * 0.52), Inches(5.2), Inches(0.4),
                     ("+ " + t), size=13, color=TEXT_MAIN)
    add_text_box(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                 "No hay plan Free. La IA va con tu propia clave (BYO): sin recargo de IA en la suscripción.",
                 size=12, italic=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 2)


def s_planes(prs):
    s = blank(prs)
    header(s, "Planes base", "Elige por tamaño — sin versión gratuita")
    rows = [
        ("Plan", "Matrículas", "Storage", "Precio", True),
        ("Starter (piloto)", "≤200", "20 GB", "$79", False),
        ("Pequeña", "≤1.000", "50 GB", "$149", False),
        ("Mediana", "≤3.000", "100 GB", "$349", False),
        ("Grande", "≤10.000", "200 GB", "$799", False),
        ("Enterprise", ">10.000", "custom", "desde $1.499", False),
    ]
    ry = Inches(2.0)
    rh = Inches(0.62)
    cols = [Inches(0.7), Inches(4.7), Inches(7.2), Inches(9.8)]
    cw = [Inches(3.8), Inches(2.3), Inches(2.3), Inches(2.8)]
    for (c0, c1, c2, c3, head) in rows:
        add_rect(s, Inches(0.5), ry, Inches(12.3), rh,
                 fill_rgb=TEAL if head else BG_SOFT, no_line=True)
        vals = [c0, c1, c2, c3]
        for k in range(4):
            add_text_box(s, cols[k], ry, cw[k], rh, vals[k], size=13,
                         bold=head or k == 0 or k == 3,
                         color=BG_WHITE if head else (TEAL_DARK if k == 3 else TEXT_MAIN),
                         anchor=MSO_ANCHOR.MIDDLE)
        ry = Emu(ry) + rh + Inches(0.08)
    add_text_box(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
                 "Modalidad Administrada (+$300/mes): ExamLab opera tu tenant. Disponible desde Pequeña.",
                 size=12, italic=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 3)


def s_bundles(prs):
    s = blank(prs)
    header(s, "Combos que tienen sentido", "Bundles recomendados por perfil de institución")
    bundles = [
        ("Colegio pequeño", "Pequeña", "$149/mes", "Sin code runner, sin aislamiento, sin SSO — el plan solo. IA con clave propia."),
        ("Facultad de Ingeniería", "Mediana + Code runner", "$398/mes", "Exámenes de programación Java/Python. SSO opcional si hay directorio."),
        ("Universidad regulada", "Grande + Aislamiento", "≈$852/mes", "Habeas Data / data residency. SSO y Reporting API ya incluidos en Grande."),
        ("Instituto con certificación", "Mediana + Certificación", "$378/mes+", "Diplomados con certificado formal + QR. IA administrada si no hay equipo TI."),
        ("Universidad grande operada", "Grande Admin + Aislamiento", "$1.198/mes", "ExamLab opera el tenant + aislamiento físico. Sin equipo propio del cliente."),
    ]
    y = Inches(1.95)
    for i, (name, combo, total, why) in enumerate(bundles):
        h = Inches(0.86)
        add_rect(s, Inches(0.5), y, Inches(12.3), h, fill_rgb=BG_SOFT, no_line=True)
        add_text_box(s, Inches(0.75), y + Inches(0.09), Inches(3.3), Inches(0.4),
                     name, size=14, bold=True, color=TEAL_DARK)
        add_text_box(s, Inches(0.75), y + Inches(0.46), Inches(3.3), Inches(0.35),
                     combo, size=11, color=TEXT_MUTED)
        add_text_box(s, Inches(4.2), y + Inches(0.24), Inches(2.0), Inches(0.4),
                     total, size=15, bold=True, color=BLUE_ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Inches(6.4), y + Inches(0.13), Inches(6.2), Inches(0.62),
                     why, size=11, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(y) + h + Inches(0.09)
    add_footer(s, 4)


def s_coherencia(prs):
    s = blank(prs)
    header(s, "Reglas de coherencia", "Qué requiere qué — y qué ya viene incluido")
    # Requiere plan mínimo
    add_rect(s, Inches(0.5), Inches(1.95), Inches(6.0), Inches(4.3),
             fill_rgb=BG_SOFT, no_line=True)
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(5.4), Inches(0.4),
                 "Add-ons que exigen un plan mínimo", size=15, bold=True, color=TEAL_DARK)
    for i, t in enumerate(["Code runner ilimitado → Mediana+",
                           "SSO / SAML → Mediana+",
                           "Aislamiento dedicado → Grande+"]):
        add_text_box(s, Inches(1.0), Inches(2.85 + i * 0.55), Inches(5.2), Inches(0.4),
                     "• " + t, size=13, color=TEXT_MAIN)
    add_text_box(s, Inches(0.8), Inches(4.7), Inches(5.4), Inches(1.4),
                 "Motivo: en planes chicos la ejecución de código ya está incluida, no hay directorio "
                 "corporativo para SSO, y el aislamiento no amortiza su setup. Se habilitan donde aportan.",
                 size=11, italic=True, color=TEXT_MUTED)
    # Incluido en tiers altos
    add_rect(s, Inches(6.83), Inches(1.95), Inches(6.0), Inches(4.3),
             fill_rgb=RGBColor(0xF5, 0xF3, 0xFF), line_rgb=TEAL_LIGHT, line_width_pt=1.5)
    add_text_box(s, Inches(7.13), Inches(2.2), Inches(5.4), Inches(0.4),
                 "Incluido sin cargo en tiers altos", size=15, bold=True, color=TEAL_DARK)
    for i, t in enumerate(["SSO / SAML → incluido en Grande",
                           "Reporting API → incluido en Grande",
                           "Onboarding guiado → 4–8h en Mediana/Grande",
                           "Backups extendidos → Mediana (30d) / Grande (90d)"]):
        add_text_box(s, Inches(7.33), Inches(2.85 + i * 0.55), Inches(5.2), Inches(0.45),
                     "✓ " + t, size=12, color=TEXT_MAIN)
    add_text_box(s, Inches(7.13), Inches(5.4), Inches(5.4), Inches(0.7),
                 "A mayor plan, más add-ons vienen incluidos: no pagas dos veces por lo mismo.",
                 size=11, italic=True, color=TEAL_DARK)
    add_footer(s, 5)


def s_cierre(prs):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=TEAL_DARK, no_line=True)
    add_text_box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.0),
                 "Armá tu plan a la medida", size=44, bold=True, color=BG_WHITE)
    add_text_box(s, Inches(1), Inches(3.9), Inches(11.3), Inches(0.8),
                 "Empezás con el plan de tu tamaño y sumás solo los add-ons que tu institución "
                 "realmente necesita. Sin versión gratuita, sin pagar por lo que no usás.",
                 size=17, color=RGBColor(0xC7, 0xD2, 0xFE))


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s_portada(prs)
    s_filosofia(prs)
    s_planes(prs)
    s_bundles(prs)
    s_coherencia(prs)
    s_cierre(prs)
    dest = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                        "..", "..", "demos", "presentacion"))
    os.makedirs(dest, exist_ok=True)
    out = os.path.join(dest, "ExamLab-Presentacion-Modelo-Modular.pptx")
    prs.save(out)
    print(f"[OK] Generado: {out}")
    print(f"     {len(prs.slides._sldIdLst)} slides · 16:9 widescreen")


if __name__ == "__main__":
    build()
