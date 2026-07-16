"""
Genera ExamLab-Presentacion-Comercial-v3.pptx con los datos del modelo
económico v3 (docs/costos/analisis/).

Uso:
    py docs/costos/analisis/_gen-presentacion.py

Requiere: python-pptx (pip install python-pptx)

Diseño: 16:9 widescreen (13.3 × 7.5 in), paleta indigo/violeta alineada al
ExamLab-Presentacion-Comercial.pptx ORIGINAL (títulos indigo oscuro #1E1B4B,
acento indigo #4F46E5, violeta #7C3AED, verde #059669), tipografía
Helvetica-like. Mismos colores que el deck original + precios/slides v3.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Paleta de marca ──────────────────────────────────────────────────
# Alineada al ExamLab-Presentacion-Comercial.pptx ORIGINAL (indigo/violeta),
# NO al teal previo. Los nombres TEAL_* se conservan por compatibilidad con el
# resto del generador, pero sus valores ya son indigo/violeta.
TEAL_DARK   = RGBColor(0x1E, 0x1B, 0x4B)   # Titular — indigo muy oscuro (original)
TEAL        = RGBColor(0x4F, 0x46, 0xE5)   # Acento primario — indigo
TEAL_LIGHT  = RGBColor(0x7C, 0x3A, 0xED)   # Acento claro / badges — violeta
BLUE_ACCENT = RGBColor(0x4F, 0x46, 0xE5)   # Precio — indigo (igual al deck original)
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT     = RGBColor(0xEE, 0xF2, 0xFF)   # Fondo suave — indigo-50
TEXT_MAIN   = RGBColor(0x1F, 0x29, 0x37)   # Texto principal — slate oscuro (original)
TEXT_MUTED  = RGBColor(0x64, 0x74, 0x8B)   # Texto atenuado — slate (original)
GOLD        = RGBColor(0xE1, 0x9A, 0x1F)   # Highlight "más popular"
GREEN_OK    = RGBColor(0x05, 0x96, 0x69)   # Checks — verde (original)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run_font(run, name="Calibri", size=Pt(14), bold=False, color=TEXT_MAIN, italic=False):
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text_box(slide, x, y, w, h, text="", *, font="Calibri", size=14, bold=False,
                 italic=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, name=font, size=Pt(size), bold=bold, italic=italic, color=color)
    return tb, tf, p


def add_rect(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None, line_width_pt=1.0, no_line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.05
    if fill_rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    if no_line:
        shape.line.fill.background()
    elif line_rgb is not None:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    shape.shadow.inherit = False
    return shape


def add_page_footer(slide, page_num, total):
    """Footer con marca + nº de página."""
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
                 "ExamLab · Plataforma educativa con IA",
                 size=9, color=TEXT_MUTED)
    add_text_box(slide, Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3),
                 f"{page_num}",
                 size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_slide_header(slide, kicker, title):
    """Kicker en teal claro + título grande en teal oscuro."""
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.4),
                 kicker.upper(),
                 size=11, bold=True, color=TEAL, font="Calibri")
    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
                 title,
                 size=28, bold=True, color=TEAL_DARK, font="Calibri")


# ── Slides individuales ─────────────────────────────────────────────

def slide_1_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Fondo con banda superior teal
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.5),
             fill_rgb=TEAL_DARK, no_line=True)
    # Título
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1),
                 "ExamLab",
                 size=54, bold=True, color=BG_WHITE, font="Calibri")
    add_text_box(slide, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.4),
                 "La plataforma educativa con IA — planes para tu institución",
                 size=16, color=RGBColor(0xC7, 0xD2, 0xFE), font="Calibri")
    # Subtítulo grande
    add_text_box(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.8),
                 "Pequeña, Mediana o Grande: hay un plan a tu medida",
                 size=32, bold=True, color=TEAL_DARK, font="Calibri")
    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                 "IA incluida (con tu API key) · Seguro y privado · Setup en días, no meses",
                 size=16, color=TEXT_MUTED, font="Calibri")
    # Etiqueta abajo
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                 "PROPUESTA COMERCIAL · Julio 2026",
                 size=11, bold=True, color=TEAL, font="Calibri")


def slide_2_por_que(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Por qué ExamLab",
                     "Tus docentes pierden horas armando y calificando")
    # Bullet points
    y = 2.1
    bullets = [
        ("Crear exámenes, talleres y proyectos consume tiempo cada semana.", TEXT_MAIN),
        ("Calificar entregas a mano es lento y poco trazable.", TEXT_MAIN),
        ("El material y el seguimiento quedan dispersos en 3-5 herramientas.", TEXT_MAIN),
        ("ExamLab automatiza con IA: generar, calificar, tutorizar y detectar copia.", TEAL_DARK),
        ("Resultado: tus docentes recuperan 8-12 horas por semana. Tu institución gana control.", TEAL_DARK),
    ]
    for txt, color in bullets:
        bold = color == TEAL_DARK
        # Bullet dot
        add_rect(slide, Inches(0.6), Inches(y + 0.13), Inches(0.12), Inches(0.12),
                 fill_rgb=TEAL_LIGHT, no_line=True)
        add_text_box(slide, Inches(0.9), Inches(y), Inches(11.5), Inches(0.5),
                     txt, size=15, bold=bold, color=color)
        y += 0.55
    # Callout box
    y_call = 5.8
    box = add_rect(slide, Inches(0.5), Inches(y_call), Inches(12.3), Inches(0.8),
                   fill_rgb=RGBColor(0xE8, 0xF6, 0xF9),
                   line_rgb=TEAL_LIGHT, line_width_pt=1.5)
    add_text_box(slide, Inches(0.8), Inches(y_call + 0.15), Inches(11.8), Inches(0.5),
                 "★  Una sola plataforma, en español, lista para tu institución — sin instalar nada.",
                 size=14, bold=True, color=TEAL_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_page_footer(slide, 1, 8)


def slide_3_plataforma(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "La plataforma",
                     "Todo lo que tu institución necesita")
    features = [
        ("📚", "Cursos, cronograma, contenidos", "Estructura por semestre + tablero visual"),
        ("📝", "Exámenes, talleres, proyectos", "Con IA para generar y calificar"),
        ("🎓", "Certificados con QR", "Verificación pública + firma digital"),
        ("🤖", "Tutor IA por curso", "Lee el material del curso y responde al alumno"),
        ("🚫", "Anti-plagio con IA", "Detección de copia entre entregas"),
        ("📊", "Libro de calificaciones", "Con pesos por corte + acta oficial exportable"),
        ("📱", "Reto en vivo", "Con PIN + QR, hasta 10k alumnos concurrentes"),
        ("✅", "Asistencia con QR rotativo", "Auto check-in del estudiante"),
        ("💬", "Mensajería + foros + broadcast", "Con etiquetas #contenido"),
        ("🖥️", "Code runner Java/Python", "Ejecución server-side en exámenes"),
    ]
    # Grid 2 columnas × 5 filas
    x_col1 = 0.5
    x_col2 = 6.9
    for i, (icon, title, sub) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(x_col1 if col == 0 else x_col2)
        y = Inches(2.1 + row * 0.9)
        # Icon
        add_text_box(slide, x, y, Inches(0.5), Inches(0.5),
                     icon, size=22)
        add_text_box(slide, x + Inches(0.55), y, Inches(5.7), Inches(0.35),
                     title, size=14, bold=True, color=TEAL_DARK)
        add_text_box(slide, x + Inches(0.55), y + Inches(0.35), Inches(5.7), Inches(0.4),
                     sub, size=11, color=TEXT_MUTED)
    add_page_footer(slide, 2, 8)


def _plan_card(slide, x, w, plan_name, target, price, matriculas,
               is_popular=False):
    """Card individual de un plan."""
    # Altura/posición ajustadas para que la card (y sus features) terminen
    # en ~6.7in, dejando aire para el footer en 6.9in (antes se solapaban).
    card_h = Inches(4.15)
    card_y = Inches(2.55)
    line_color = GOLD if is_popular else RGBColor(0xD0, 0xD9, 0xE1)
    fill_color = BG_WHITE if not is_popular else RGBColor(0xFF, 0xFB, 0xF0)
    add_rect(slide, x, card_y, w, card_h,
             fill_rgb=fill_color,
             line_rgb=line_color, line_width_pt=2.0 if is_popular else 1.0)
    # Badge MÁS POPULAR
    if is_popular:
        badge_w = Inches(1.5)
        badge = add_rect(slide, x + (w - badge_w) / 2, card_y - Inches(0.2),
                         badge_w, Inches(0.35),
                         fill_rgb=GOLD, no_line=True)
        add_text_box(slide, x + (w - badge_w) / 2, card_y - Inches(0.18),
                     badge_w, Inches(0.3),
                     "MÁS POPULAR",
                     size=9, bold=True, color=BG_WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Nombre plan
    add_text_box(slide, x, card_y + Inches(0.30), w, Inches(0.5),
                 plan_name,
                 size=26, bold=True, color=TEAL_DARK,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, x, card_y + Inches(0.85), w, Inches(0.35),
                 target,
                 size=12, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
    # Precio
    add_text_box(slide, x, card_y + Inches(1.25), w, Inches(0.8),
                 price,
                 size=44, bold=True, color=BLUE_ACCENT,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, x, card_y + Inches(2.05), w, Inches(0.3),
                 "USD / mes",
                 size=12, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
    # Divisor
    div = slide.shapes.add_connector(1, x + Inches(0.4), card_y + Inches(2.45),
                                     x + w - Inches(0.4), card_y + Inches(2.45))
    div.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEC)
    # Matrículas
    add_text_box(slide, x, card_y + Inches(2.55), w, Inches(0.3),
                 "Matrículas activas",
                 size=10, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, x, card_y + Inches(2.85), w, Inches(0.4),
                 matriculas,
                 size=20, bold=True, color=TEAL_DARK,
                 align=PP_ALIGN.CENTER)
    # Features check
    y_check = card_y + Inches(3.35)
    for feat in ["Todas las funciones incluidas",
                 "Cursos y usuarios ilimitados",
                 "Soporte + backups incluidos"]:
        add_text_box(slide, x + Inches(0.2), y_check, Inches(0.3), Inches(0.25),
                     "✓", size=12, bold=True, color=GREEN_OK)
        add_text_box(slide, x + Inches(0.5), y_check, w - Inches(0.5), Inches(0.25),
                     feat, size=10, color=TEXT_MAIN)
        y_check += Inches(0.25)


def slide_4_planes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Planes",
                     "Un precio para cada tamaño de institución")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.4),
                 "Mismas funciones en todos los planes — solo cambia el tope de matrículas.",
                 size=13, color=TEXT_MUTED)
    # 3 cards
    card_w = Inches(3.95)
    gap = Inches(0.25)
    x1 = Inches(0.5)
    x2 = x1 + card_w + gap
    x3 = x2 + card_w + gap
    _plan_card(slide, x1, card_w, "Pequeña", "Institución pequeña",
               "$149", "Hasta 1.000")
    _plan_card(slide, x2, card_w, "Mediana", "Institución mediana",
               "$349", "Hasta 3.000", is_popular=True)
    _plan_card(slide, x3, card_w, "Grande", "Institución grande",
               "$799", "Hasta 10.000")
    # Footer Enterprise — debajo de las cards (que terminan en ~6.7in), sin solaparse.
    add_text_box(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
                 "¿Más de 10.000 matrículas o requerimientos regulatorios? Plan Enterprise custom desde $1.499/mes · Contáctanos →",
                 size=10, italic=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)


def slide_5_todo_incluido(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Todo incluido",
                     "Todas las funciones, en todos los planes")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.4),
                 "Sin funciones bloqueadas: la única diferencia entre planes es el tope de matrículas.",
                 size=13, color=TEXT_MUTED)
    features = [
        "Cursos, cronograma y tablero",
        "Exámenes, talleres y proyectos",
        "Generación de evaluaciones con IA",
        "Generación de contenido (PPTX/guías) con IA",
        "Calificación automática con retroalimentación",
        "Tutor IA del curso (lee el material)",
        "Detección de copia / antifraude",
        "Banco de preguntas reutilizable",
        "Asistencia con QR rotativo (auto check-in)",
        "Sincronización con Google/Microsoft Calendar",
        "Mensajería 1-a-1, foros y difusión",
        "Encuestas + Reto en vivo con PIN",
        "Certificados con QR verificable",
        "Ejecución de código (Java/Python) en línea",
        "Multi-sede, branding y auditoría",
        "Notificaciones in-app + push + email",
    ]
    # Grid 4 columnas × 4 filas
    col_w = Inches(3.05)
    for i, feat in enumerate(features):
        col = i % 4
        row = i // 4
        x = Inches(0.5 + col * 3.2)
        y = Inches(2.3 + row * 0.55)
        add_text_box(slide, x, y, Inches(0.3), Inches(0.35),
                     "✓", size=14, bold=True, color=GREEN_OK)
        add_text_box(slide, x + Inches(0.35), y, col_w - Inches(0.35), Inches(0.35),
                     feat, size=11, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)
    # Callout final
    y_call = 5.7
    box = add_rect(slide, Inches(0.5), Inches(y_call), Inches(12.3), Inches(1.0),
                   fill_rgb=RGBColor(0xE8, 0xF6, 0xF9),
                   line_rgb=TEAL_LIGHT, line_width_pt=1.5)
    add_text_box(slide, Inches(0.8), Inches(y_call + 0.15), Inches(11.8), Inches(0.35),
                 "★  BYO API Key (Bring Your Own Key)",
                 size=13, bold=True, color=TEAL_DARK)
    add_text_box(slide, Inches(0.8), Inches(y_call + 0.5), Inches(11.8), Inches(0.4),
                 "La institución conecta su propia clave de Gemini/OpenAI. Sin sobrecosto de IA en la suscripción.",
                 size=11, color=TEXT_MAIN)
    add_page_footer(slide, 4, 8)


def slide_6_ia_flexible(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Modelo de IA",
                     "Tú decides cuánto se usa la IA — y cuánto pagas")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5),
                 "Dos maneras de usar la IA, sin sorpresas de costo:",
                 size=14, bold=True, color=TEAL_DARK)
    # 2 cards lado a lado
    card_y = Inches(2.4)
    card_h = Inches(3.5)
    card_w = Inches(6.0)
    # Card 1: BYO
    add_rect(slide, Inches(0.5), card_y, card_w, card_h,
             fill_rgb=RGBColor(0xF0, 0xF9, 0xFB),
             line_rgb=TEAL, line_width_pt=2.0)
    add_text_box(slide, Inches(0.8), card_y + Inches(0.25), card_w - Inches(0.4),
                 Inches(0.5),
                 "🔑  BYO API Key (recomendado)",
                 size=18, bold=True, color=TEAL_DARK)
    add_text_box(slide, Inches(0.8), card_y + Inches(0.85), card_w - Inches(0.4),
                 Inches(0.4),
                 "Conectas tu propia clave de Gemini/OpenAI.",
                 size=13, color=TEXT_MAIN)
    bullets_byo = [
        "Sobre-costo de IA en la suscripción: $0",
        "Pagas al proveedor directo (Google/OpenAI) según uso real",
        "Costo típico ~$0.06/matrícula/mes con Gemini Flash",
        "Cola sync/async para amortiguar picos de gasto",
        "Failover automático a claves de respaldo",
    ]
    y_b = card_y + Inches(1.4)
    for b in bullets_byo:
        add_text_box(slide, Inches(0.9), y_b, Inches(0.3), Inches(0.3),
                     "✓", size=12, bold=True, color=GREEN_OK)
        add_text_box(slide, Inches(1.25), y_b, card_w - Inches(0.9),
                     Inches(0.35),
                     b, size=11, color=TEXT_MAIN)
        y_b += Inches(0.4)
    # Card 2: IA administrada
    x2 = Inches(6.83)
    add_rect(slide, x2, card_y, card_w, card_h,
             fill_rgb=BG_WHITE,
             line_rgb=RGBColor(0xD0, 0xD9, 0xE1), line_width_pt=1.0)
    add_text_box(slide, x2 + Inches(0.3), card_y + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.5),
                 "🤖  IA administrada (add-on)",
                 size=18, bold=True, color=TEAL_DARK)
    add_text_box(slide, x2 + Inches(0.3), card_y + Inches(0.85),
                 card_w - Inches(0.4), Inches(0.4),
                 "ExamLab gestiona la clave. Facturación medida por matrícula.",
                 size=13, color=TEXT_MAIN)
    bullets_ad = [
        "$0.10 por matrícula activa / mes",
        "Sin gestión de Google Cloud del lado del cliente",
        "Tope de consumo configurable por institución",
        "Facturación incluida en la factura de ExamLab",
        "Ideal para instituciones sin equipo TI dedicado",
    ]
    y_b = card_y + Inches(1.4)
    for b in bullets_ad:
        add_text_box(slide, x2 + Inches(0.4), y_b, Inches(0.3), Inches(0.3),
                     "✓", size=12, bold=True, color=GREEN_OK)
        add_text_box(slide, x2 + Inches(0.75), y_b, card_w - Inches(0.9),
                     Inches(0.35),
                     b, size=11, color=TEXT_MAIN)
        y_b += Inches(0.4)
    # Cierre
    add_text_box(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.4),
                 "Cola sync/async · Failover multi-key · Auditoría completa · Datos nunca dejan tu institución",
                 size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_page_footer(slide, 5, 8)


def slide_7_comparativa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Comparativa",
                     "Más barato que Moodle Cloud. Mucho más barato que Canvas.")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.4),
                 "Precio USD/matrícula/mes para una institución con 3.000 matrículas.",
                 size=13, color=TEXT_MUTED)
    # Tabla comparativa
    rows = [
        ("Producto", "USD/mes total", "Por matrícula/mes", "Notas"),
        ("Chamilo self-hosted", "~$350 (valorizado)", "$0.12", "Gratis + admin server + VPS"),
        ("ExamLab Mediana", "$349", "$0.12", "IA + Reto en vivo + antifraude incluidos"),
        ("Moodle Cloud Standard (750)", "$173", "$0.23", "Sin IA nativa, plugin extra"),
        ("Canvas Small (negociado)", "$2.500", "$0.83", "Contract 3 años típico"),
        ("Blackboard mid-market", "$5.000+", "$1.67+", "Setup + capacitación aparte"),
    ]
    y = 2.4
    row_h = 0.55
    for i, row in enumerate(rows):
        header = i == 0
        highlight = row[0] == "ExamLab Mediana"
        bg = TEAL_DARK if header else (RGBColor(0xE8, 0xF6, 0xF9) if highlight else (
            BG_WHITE if i % 2 == 1 else BG_SOFT))
        text_color = BG_WHITE if header else (TEAL_DARK if highlight else TEXT_MAIN)
        bold = header or highlight
        add_rect(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(row_h),
                 fill_rgb=bg, no_line=True)
        # 4 columnas
        widths = [Inches(4.5), Inches(2.8), Inches(2.5), Inches(2.5)]
        cx = Inches(0.7)
        for col_i, (val, w) in enumerate(zip(row, widths)):
            align = PP_ALIGN.LEFT if col_i == 0 else (
                PP_ALIGN.CENTER if col_i < 3 else PP_ALIGN.LEFT)
            add_text_box(slide, cx, Inches(y), w, Inches(row_h),
                         val, size=12, bold=bold, color=text_color,
                         align=align, anchor=MSO_ANCHOR.MIDDLE)
            cx += w
        y += row_h
    # Fuente
    add_text_box(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                 "Fuentes: moodlecloud.com/standard-plans · vendr.com/marketplace/canvas · benchmarks públicos 2026-07",
                 size=9, italic=True, color=TEXT_MUTED)
    add_page_footer(slide, 6, 8)


def slide_8_valor_ahorro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Cuánto ahorras",
                     "El costo real de NO tener ExamLab")
    # Row 1: métricas
    row_y = Inches(2.1)
    metrics = [
        ("$60k+", "Canvas negociado / año\npara 5.000 alumnos"),
        ("8-12h", "Que un docente ahorra\npor semana con IA"),
        ("$0", "Costo de IA en el plan\ncon BYO API key"),
        ("1 día", "Setup completo con\nasistencia guiada"),
    ]
    card_w = Inches(3.0)
    gap = Inches(0.1)
    for i, (num, label) in enumerate(metrics):
        x = Inches(0.5 + i * 3.1)
        add_rect(slide, x, row_y, card_w, Inches(1.7),
                 fill_rgb=RGBColor(0xEE, 0xF2, 0xFF), no_line=True)
        add_text_box(slide, x, row_y + Inches(0.15), card_w, Inches(0.6),
                     num, size=32, bold=True, color=BLUE_ACCENT,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, x, row_y + Inches(0.9), card_w, Inches(0.7),
                     label, size=11, color=TEXT_MAIN,
                     align=PP_ALIGN.CENTER)
    # Ejemplo real
    box_y = Inches(4.3)
    add_rect(slide, Inches(0.5), box_y, Inches(12.3), Inches(2.3),
             fill_rgb=RGBColor(0xF9, 0xFA, 0xFB),
             line_rgb=TEAL_LIGHT, line_width_pt=1.5)
    add_text_box(slide, Inches(0.8), box_y + Inches(0.2), Inches(11.8),
                 Inches(0.4),
                 "Ejemplo: Universidad regional con 3.000 matrículas activas",
                 size=15, bold=True, color=TEAL_DARK)
    lines = [
        ("Plan Mediana Auto (12 meses):", "$349 × 12 = $4.188/año"),
        ("Comparable Canvas negociado:", "$24.000-$30.000/año"),
        ("Ahorro anual con ExamLab:", "~$20.000-$25.000 USD"),
        ("Setup + migración:", "Incluidos (Admin) vs 3-6 meses Canvas"),
    ]
    y_l = box_y + Inches(0.75)
    for label, val in lines:
        add_text_box(slide, Inches(0.9), y_l, Inches(6), Inches(0.3),
                     label, size=12, color=TEXT_MAIN)
        add_text_box(slide, Inches(7), y_l, Inches(5.5), Inches(0.3),
                     val, size=12, bold=True, color=TEAL_DARK)
        y_l += Inches(0.35)
    add_page_footer(slide, 7, 8)


def slide_9_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Banda superior
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.5),
             fill_rgb=TEAL_DARK, no_line=True)
    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
                 "Dale tiempo a tus docentes.",
                 size=40, bold=True, color=BG_WHITE)
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7),
                 "Control a tu institución.",
                 size=40, bold=True, color=RGBColor(0xC7, 0xD2, 0xFE))
    # Bloque central
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.5),
                 "Empieza con el plan que se ajusta a tu tamaño y crece cuando lo necesites.",
                 size=18, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
    # 3 puntos clave
    y = 3.9
    puntos = [
        ("★", "IA para generar, calificar, tutorizar y detectar copia"),
        ("★", "Seguridad de nivel empresarial y datos privados por institución"),
        ("★", "Setup en días · Soporte en español · Sin contract mínimo"),
    ]
    for icon, txt in puntos:
        add_text_box(slide, Inches(2.5), Inches(y), Inches(0.4), Inches(0.4),
                     icon, size=18, color=TEAL_LIGHT,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3.0), Inches(y), Inches(9.5), Inches(0.4),
                     txt, size=14, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.55
    # CTA
    cta_y = Inches(6.0)
    add_rect(slide, Inches(3.5), cta_y, Inches(6.3), Inches(0.8),
             fill_rgb=TEAL, no_line=True)
    add_text_box(slide, Inches(3.5), cta_y + Inches(0.15), Inches(6.3),
                 Inches(0.5),
                 "Solicita una demo · demo@examlab.co",
                 size=18, bold=True, color=BG_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
                 "ExamLab · Plataforma educativa con IA · Julio 2026",
                 size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────

def slide_storage(prs):
    """Almacenamiento incluido por plan + precio de storage extra (cliente)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Almacenamiento",
                     "Cuánto espacio incluye cada plan")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.45),
                 "El material y las entregas se guardan en la plataforma. Los videos van por enlace "
                 "externo (YouTube/Drive) — no ocupan tu espacio.",
                 size=13, color=TEXT_MUTED)
    # 3 cards de storage por plan
    cards = [
        ("Pequeña", "50 GB", "~7× lo que usa una institución de 1.000 matrículas", TEAL),
        ("Mediana", "100 GB", "~5× lo que usa una institución de 3.000 matrículas", BLUE_ACCENT),
        ("Grande", "200 GB", "~3× lo que usa una institución de 10.000 matrículas", TEAL_DARK),
    ]
    card_w = Inches(3.95)
    gap = Inches(0.25)
    xs = [Inches(0.5), Inches(0.5) + card_w + gap, Inches(0.5) + 2 * (card_w + gap)]
    card_y = Inches(2.45)
    card_h = Inches(2.35)
    for (name, gb, note, accent), x in zip(cards, xs):
        add_rect(slide, x, card_y, card_w, card_h, fill_rgb=BG_SOFT,
                 line_rgb=accent, line_width_pt=1.5)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.25), card_w - Inches(0.6),
                     Inches(0.4), name, size=15, bold=True, color=TEXT_MUTED)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.7), card_w - Inches(0.6),
                     Inches(0.7), gb, size=40, bold=True, color=accent)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.55), card_w - Inches(0.6),
                     Inches(0.7), note, size=11, color=TEXT_MAIN)
    # Callout: storage extra + tip video externo
    y_call = 5.25
    add_rect(slide, Inches(0.5), Inches(y_call), Inches(12.3), Inches(1.4),
             fill_rgb=RGBColor(0xE8, 0xF6, 0xF9), line_rgb=TEAL_LIGHT, line_width_pt=1.5)
    add_text_box(slide, Inches(0.8), Inches(y_call + 0.18), Inches(11.7), Inches(0.4),
                 "¿Necesitas más espacio?  Storage adicional: $10 / 100 GB al mes.",
                 size=15, bold=True, color=TEAL_DARK)
    add_text_box(slide, Inches(0.8), Inches(y_call + 0.62), Inches(11.7), Inches(0.7),
                 "El espacio incluido cubre con holgura a más del 99% de las instituciones. "
                 "Antes de ampliar, siempre conviene enlazar los videos y el material pesado por URL "
                 "externa (gratis). El add-on de storage es solo para material propietario que no se "
                 "pueda externalizar. Enterprise: almacenamiento a medida.",
                 size=11, color=TEXT_MAIN)
    add_page_footer(slide, 5, 10)


def slide_ia_costo(prs):
    """Costo de IA (BYO) — hero $0 recargo + tabla de estimado por escala."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "IA, sin letra chica",
                     "El costo de IA es bajo, transparente y bajo tu control")
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.55),
                 "Con el modelo BYO, tu suscripción ExamLab no incluye ningún recargo por IA: "
                 "conectas tu propia clave de Gemini/OpenAI y pagas solo el uso real, directo al proveedor.",
                 size=13, color=TEXT_MUTED)
    # Hero $0
    hero_y = Inches(2.45)
    add_rect(slide, Inches(0.5), hero_y, Inches(4.1), Inches(3.1),
             fill_rgb=RGBColor(0xF0, 0xF9, 0xFB), line_rgb=TEAL, line_width_pt=2.0)
    add_text_box(slide, Inches(0.7), hero_y + Inches(0.35), Inches(3.7), Inches(0.5),
                 "Recargo de IA en tu plan", size=14, bold=True, color=TEAL_DARK,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7), hero_y + Inches(0.95), Inches(3.7), Inches(1.2),
                 "$0", size=90, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(0.7), hero_y + Inches(2.35), Inches(3.7), Inches(0.6),
                 "Pagas la IA directo a Google/OpenAI, según tu consumo real.",
                 size=11, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
    # Tabla estimado (col der)
    tx = Inches(4.95)
    tw = Inches(7.85)
    add_text_box(slide, tx, hero_y - Inches(0.05), tw, Inches(0.4),
                 "Estimado de costo de IA al mes (modo BYO)", size=13, bold=True, color=TEAL_DARK)
    rows = [
        ("Matrículas activas", "Uso típico", "Uso intensivo", True),
        ("200", "~$12", "~$40", False),
        ("1.000", "~$60", "~$200", False),
        ("3.000", "~$180", "~$600", False),
        ("10.000", "~$600", "~$2.000", False),
    ]
    ry = hero_y + Inches(0.45)
    rh = Inches(0.5)
    col_x = [tx + Inches(0.15), tx + Inches(3.5), tx + Inches(5.7)]
    for (c0, c1, c2, is_head) in rows:
        add_rect(slide, tx, ry, tw, rh,
                 fill_rgb=TEAL if is_head else BG_SOFT, no_line=True)
        hc = BG_WHITE if is_head else TEXT_MAIN
        add_text_box(slide, col_x[0], ry, Inches(3.3), rh, c0, size=12,
                     bold=is_head, color=hc, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, col_x[1], ry, Inches(2.1), rh, c1, size=12,
                     bold=(is_head or True) if not is_head else True,
                     color=(BG_WHITE if is_head else TEAL_DARK), anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, col_x[2], ry, Inches(2.0), rh, c2, size=12,
                     color=(BG_WHITE if is_head else TEXT_MUTED), anchor=MSO_ANCHOR.MIDDLE)
        ry += rh + Inches(0.06)
    add_text_box(slide, tx, ry + Inches(0.05), tw, Inches(0.75),
                 "Valores estimados en USD/mes, pagados directamente a tu proveedor de IA (Google/OpenAI) "
                 "según tu consumo real — no son un cargo de ExamLab. ¿Prefieres no administrar la clave? "
                 "Con el add-on de IA administrada, ExamLab se encarga de todo por $0.10 por matrícula/mes.",
                 size=10, italic=True, color=TEXT_MUTED)
    add_page_footer(slide, 7, 11)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_portada(prs)
    slide_2_por_que(prs)
    slide_3_plataforma(prs)
    slide_4_planes(prs)
    slide_storage(prs)          # storage incluido por plan + extra (cliente)
    slide_5_todo_incluido(prs)
    slide_6_ia_flexible(prs)
    slide_ia_costo(prs)         # costo de IA (BYO) — $0 recargo + estimado por escala
    slide_7_comparativa(prs)
    slide_8_valor_ahorro(prs)
    slide_9_cierre(prs)

    # Salida: la presentación comercial vive en docs/demos/presentacion (junto a
    # las demás presentaciones). El generador y las fuentes de datos quedan en
    # docs/costos/analisis.
    dest = os.path.abspath(os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "..", "..", "demos", "presentacion"))
    os.makedirs(dest, exist_ok=True)
    out = os.path.join(dest, "ExamLab-Presentacion-Comercial-v3.pptx")
    prs.save(out)
    print(f"[OK] Generado: {out}")
    print(f"     {len(prs.slides)} slides · 16:9 widescreen")


if __name__ == "__main__":
    build()
